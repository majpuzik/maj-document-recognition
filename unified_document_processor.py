#!/usr/bin/env python3
"""
NAS5 Docker Apps Collection
"""

"""
This project implements a Model Context Protocol (MCP) server that allows interaction with Gmail accounts via IMAP and SMTP. It provides tools for searching emails, retrieving content, managing labels
"""

"""
This project implements a Model Context Protocol (MCP) server that allows interaction with Gmail accounts via IMAP and SMTP. It provides tools for searching emails, retrieving content, managing labels
"""

"""
Unified Document Processor - Paperless-NGX Style
================================================
Single workflow for ALL document sources:
- Email attachments (mbox, thunderbird)
- Local files (OneDrive, Dropbox, local folders)
- Scanned documents
- Any PDF, Office, Image file

Features:
- 3 local LLM consensus classification
- Recipient extraction (Pan/Pani/Firma)
- Email metadata extraction
- OCR for images and scanned PDFs
- Paperless-ngx integration
- Deduplication

Version: 3.0.0
"""

import sys
import json
import hashlib
import logging
import subprocess
import platform
import os
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any, Union
from dataclasses import dataclass, asdict, field
from concurrent.futures import ThreadPoolExecutor
import requests


# =============================================================================
# PLATFORM DETECTION & OPTIMIZATION
# =============================================================================

def detect_platform() -> str:
    """
    Detect platform for optimal LLM backend selection.

    Returns:
        "mac_apple_silicon" - Use MLX or llama.cpp with Metal
        "linux_nvidia" - Use vLLM for batch processing
        "generic" - Use standard Ollama
    """
    system = platform.system()
    machine = platform.machine()

    if system == "Darwin":
        if machine == "arm64":
            return "mac_apple_silicon"
        return "mac_intel"
    elif system == "Linux":
        # Check for NVIDIA GPU
        try:
            result = subprocess.run(
                ["nvidia-smi", "--query-gpu=name", "--format=csv,noheader"],
                capture_output=True, text=True, timeout=5
            )
            if result.returncode == 0 and result.stdout.strip():
                return "linux_nvidia"
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass
        return "linux_cpu"
    return "generic"


def get_optimal_settings(platform_type: str) -> Dict:
    """
    Get optimal LLM settings for the detected platform.

    Mac Apple Silicon:
        - Use smaller models (8B max recommended)
        - Enable Metal acceleration (OLLAMA_METAL=1)
        - Use Q4_K_M quantization for speed
        - Optional: Use MLX for fastest inference

    Linux NVIDIA (DGX):
        - Can use larger models (32B+)
        - Use vLLM for 3-5x speedup on batch processing
        - GPU memory pooling enabled

    Generic:
        - Standard Ollama settings
        - Conservative model sizes
    """
    settings = {
        "platform": platform_type,
        "recommended_models": [],
        "env_vars": {},
        "use_vllm": False,
        "use_mlx": False,
        "max_model_size_gb": 8,
        "batch_size": 1,
        "parallel_requests": 1,
    }

    if platform_type == "mac_apple_silicon":
        settings.update({
            "recommended_models": [
                "ministral-3:8b",       # 6GB, fast
                "llama3.2:3b",          # 2GB, ultra fast
                "qwen2.5:7b",           # 5GB
            ],
            "env_vars": {
                "OLLAMA_METAL": "1",           # Enable Metal GPU
                "OLLAMA_NUM_PARALLEL": "2",    # Conservative parallelism
                "OLLAMA_MAX_LOADED_MODELS": "1", # One model at a time
            },
            "use_mlx": True,  # Can use MLX if installed
            "max_model_size_gb": 10,  # M1/M2/M3 unified memory
            "parallel_requests": 2,
        })
    elif platform_type == "linux_nvidia":
        settings.update({
            "recommended_models": [
                "qwen2.5:32b",          # High quality
                "ministral-3:8b",       # Fast
                "czech-finance-speed",  # Czech specialized
            ],
            "env_vars": {
                "OLLAMA_NUM_PARALLEL": "4",
                "OLLAMA_MAX_LOADED_MODELS": "2",
                "CUDA_VISIBLE_DEVICES": "0",  # Use first GPU
            },
            "use_vllm": True,  # vLLM for batch processing
            "max_model_size_gb": 48,  # A100/H100
            "batch_size": 8,
            "parallel_requests": 4,
        })
    else:
        settings.update({
            "recommended_models": [
                "llama3.2:3b",
                "ministral-3:3b",
            ],
            "max_model_size_gb": 4,
        })

    return settings


# Detect platform at module load
PLATFORM_TYPE = detect_platform()
PLATFORM_SETTINGS = get_optimal_settings(PLATFORM_TYPE)

# Setup logging EARLY (before MLX section needs it)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


# =============================================================================
# MLX BACKEND (Apple Silicon Only - Fastest)
# =============================================================================

# MLX model paths - searched in order
MLX_MODEL_PATHS = [
    Path.home() / "mlx-models" / "qwen2.5-3b-4bit",
    Path.home() / "mlx-models" / "llama-3.2-3b-4bit",
    Path.home() / "mlx-models" / "ministral-8b-4bit",
    Path("/opt/mlx-models"),  # Shared location
]

def detect_mlx_model() -> Optional[Path]:
    """
    Detect if MLX model is available (Mac Apple Silicon only).
    Returns path to first available MLX model or None.
    """
    if PLATFORM_TYPE != "mac_apple_silicon":
        return None

    try:
        import mlx_lm
    except ImportError:
        return None

    for model_path in MLX_MODEL_PATHS:
        if model_path.exists() and (model_path / "config.json").exists():
            return model_path

    return None


class MLXClassifier:
    """
    MLX-based classifier for Apple Silicon.
    2-3x faster than Ollama on Mac M1/M2/M3/M4.
    """

    def __init__(self, model_path: Path):
        self.model_path = model_path
        self.model = None
        self.tokenizer = None
        self._load_model()

    def _load_model(self):
        """Load MLX model and tokenizer"""
        try:
            from mlx_lm import load
            self.model, self.tokenizer = load(str(self.model_path))
            logger.info(f"✅ MLX model loaded: {self.model_path.name}")
        except Exception as e:
            logger.error(f"❌ Failed to load MLX model: {e}")
            raise

    def generate(self, prompt: str, max_tokens: int = 500) -> str:
        """Generate response using MLX model"""
        try:
            from mlx_lm import generate
            response = generate(
                self.model,
                self.tokenizer,
                prompt=prompt,
                max_tokens=max_tokens,
            )
            return response.strip()
        except Exception as e:
            logger.warning(f"MLX generation failed: {e}")
            return ""


# Auto-detect MLX model at module load
MLX_MODEL_PATH = detect_mlx_model()
MLX_AVAILABLE = MLX_MODEL_PATH is not None

if MLX_AVAILABLE:
    logger.info(f"🚀 MLX model detected: {MLX_MODEL_PATH.name} (fastest inference)")
else:
    if PLATFORM_TYPE == "mac_apple_silicon":
        logger.info("💡 MLX not available. Install model: python3 -m mlx_lm.convert --hf-path Qwen/Qwen2.5-3B-Instruct -q --mlx-path ~/mlx-models/qwen2.5-3b-4bit")


# =============================================================================
# CONFIGURATION
# =============================================================================

# Supported file types
SUPPORTED_DOCUMENTS = {
    # PDF
    '.pdf': 'pdf',
    # Word
    '.doc': 'word', '.docx': 'word',
    # Excel
    '.xls': 'excel', '.xlsx': 'excel',
    # PowerPoint
    '.ppt': 'ppt', '.pptx': 'ppt',
    # Images (OCR)
    '.jpg': 'image', '.jpeg': 'image', '.png': 'image',
    '.tiff': 'image', '.tif': 'image', '.bmp': 'image',
    # Text
    '.txt': 'text', '.rtf': 'text', '.csv': 'text',
    # OpenDocument
    '.odt': 'odt', '.ods': 'ods',
}

# Local Ollama models - PLATFORM-AWARE CONFIGURATION
# Models are selected based on detected platform (Mac vs DGX)
if PLATFORM_TYPE == "mac_apple_silicon":
    # Mac M1/M2/M3/M4: Use smaller models, Metal acceleration
    OLLAMA_MODELS = [
        "ministral-3:8b",    # 6GB, fast with Metal
        "llama3.2:3b",       # 2GB, ultra fast
        "qwen2.5:7b",        # 5GB, good quality
    ]
    FALLBACK_MODELS = [
        "llama3.1:8b",
        "phi3:mini",
    ]
    FAST_MODE = True  # Always fast on Mac (single model)
elif PLATFORM_TYPE == "linux_nvidia":
    # DGX/Linux with NVIDIA GPU: Can use larger models
    OLLAMA_MODELS = [
        "ministral-3:8b",              # Fast, good quality
        "qwen2.5:14b",                 # Better quality
        "czech-finance-speed:latest",  # Czech specialist
    ]
    FALLBACK_MODELS = [
        "qwen2.5:32b",     # High quality if GPU memory allows
        "llama3.1:8b",
    ]
    FAST_MODE = True  # Fast mode due to parallel scans
else:
    # Generic/CPU: Smallest models only
    OLLAMA_MODELS = [
        "llama3.2:3b",
        "ministral-3:3b",
    ]
    FALLBACK_MODELS = []
    FAST_MODE = True

# Log platform configuration at startup
logger.info(f"🖥️  Platform: {PLATFORM_TYPE}")
logger.info(f"📊 Models: {OLLAMA_MODELS}")
logger.info(f"⚡ FAST_MODE: {FAST_MODE}")

# =============================================================================
# HIERARCHICAL MODEL ORCHESTRATOR ("Dirigent")
# =============================================================================
# Strategy: Fast model first → escalate to big model only if low confidence
# Saves ~50% inference time while maintaining accuracy

HIERARCHICAL_MODE = os.environ.get('HIERARCHICAL_MODE', '1') == '1'
ESCALATION_THRESHOLD = float(os.environ.get('ESCALATION_THRESHOLD', '0.7'))  # Escalate if confidence < 70%

# Fast models (2-3x faster, ~90% accuracy for clear documents)
FAST_MODELS_HIERARCHY = [
    "qwen2.5:3b-instruct-q4_K_M",         # 3B quantized - fastest with good Czech
    "ministral:3b-instruct-2410-q4_K_M",  # 3B Mistral quantized
    "llama3.2:3b-instruct-q4_K_M",        # 3B Llama quantized
    "gemma2:2b-instruct-q4_K_M",          # 2B - very fast
    "qwen2.5:3b",                         # 3B default
    "llama3.2:3b",                        # 3B default
]

# Big models (escalation targets - higher accuracy for ambiguous documents)
BIG_MODELS_HIERARCHY = [
    "qwen2.5:14b-instruct-q4_K_M",        # 14B - best accuracy
    "qwen2.5:32b-instruct-q4_K_M",        # 32B - highest accuracy (DGX only)
    "mistral:7b-instruct-q4_K_M",         # 7B - good balance
    "llama3.1:8b",                        # 8B default
    "qwen2.5:7b",                         # 7B - good Czech
]

logger.info(f"🎭 HIERARCHICAL_MODE: {HIERARCHICAL_MODE} (threshold: {ESCALATION_THRESHOLD})")

# Document types
DOC_TYPES = [
    "INVOICE", "RECEIPT", "CONTRACT", "INSURANCE", "MEDICAL_REPORT",
    "BANK_STATEMENT", "TAX_DOCUMENT", "LEGAL_DOCUMENT", "CORRESPONDENCE",
    "TECHNICAL_MANUAL", "MARKETING", "ORDER", "DELIVERY_NOTE", "WARRANTY",
    "CERTIFICATE", "APPLICATION_FORM", "REPORT", "MINUTES", "PERSONAL_ID",
    "POWER_OF_ATTORNEY", "DECLARATION", "NOTIFICATION", "EMAIL", "OTHER"
]

# Recipient type mapping for Paperless tags
RECIPIENT_TYPES = {
    "PERSON_MALE": "Pan",
    "PERSON_FEMALE": "Paní",
    "COMPANY": "Firma",
    "INSTITUTION": "Instituce",
    "UNKNOWN": "Neznámý"
}


# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class DocumentSource:
    """Source information for the document"""
    type: str  # "email", "file", "scan"
    path: str
    # Email-specific
    email_sender: Optional[str] = None
    email_recipient: Optional[str] = None
    email_subject: Optional[str] = None
    email_date: Optional[str] = None
    # File-specific
    original_filename: Optional[str] = None
    parent_folder: Optional[str] = None


@dataclass
class RecipientInfo:
    """Extracted recipient information"""
    name: str
    type: str  # PERSON_MALE, PERSON_FEMALE, COMPANY, INSTITUTION
    confidence: float
    source: str  # "email", "content", "filename"


@dataclass
class ClassificationResult:
    """Complete classification result"""
    # Identifiers
    file_path: str
    file_hash: str
    file_type: str

    # Source
    source: DocumentSource

    # Classification
    document_type: str
    type_confidence: float
    type_votes: Dict[str, str] = field(default_factory=dict)
    type_consensus: bool = False

    # Recipient
    recipient: Optional[RecipientInfo] = None

    # Suggested metadata
    suggested_title: Optional[str] = None
    suggested_tags: List[str] = field(default_factory=list)

    # Extracted text (preview)
    text_preview: str = ""
    text_length: int = 0

    # Processing info
    processed_at: str = ""
    processing_time_ms: int = 0
    ocr_used: bool = False
    models_used: List[str] = field(default_factory=list)


# =============================================================================
# TEXT EXTRACTION
# =============================================================================

class UnifiedTextExtractor:
    """Extract text from any supported document type"""

    def __init__(self, use_ocr: bool = True):
        self.use_ocr = use_ocr
        self.ocr_languages = "ces+eng+deu"

    def extract(self, file_path: Path) -> Dict[str, Any]:
        """
        Extract text from any supported file

        Returns:
            {
                'text': str,
                'success': bool,
                'method': str,  # 'native', 'ocr', 'direct'
                'confidence': float (for OCR)
            }
        """
        ext = file_path.suffix.lower()
        file_type = SUPPORTED_DOCUMENTS.get(ext)

        if not file_type:
            return {'text': '', 'success': False, 'method': 'unsupported'}

        try:
            if file_type == 'pdf':
                return self._extract_pdf(file_path)
            elif file_type == 'word':
                return self._extract_word(file_path)
            elif file_type == 'excel':
                return self._extract_excel(file_path)
            elif file_type == 'image':
                return self._extract_image(file_path)
            elif file_type == 'text':
                return self._extract_text(file_path)
            elif file_type in ('odt', 'ods'):
                return self._extract_odf(file_path)
            elif file_type == 'ppt':
                return self._extract_ppt(file_path)
            else:
                return {'text': '', 'success': False, 'method': 'unknown'}
        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {e}")
            return {'text': '', 'success': False, 'method': 'error', 'error': str(e)}

    def _extract_pdf(self, file_path: Path) -> Dict:
        """Extract text from PDF (native first, then OCR)"""
        # Try pdftotext first (fast, native text)
        try:
            result = subprocess.run(
                ['pdftotext', '-layout', str(file_path), '-'],
                capture_output=True, text=True, timeout=60
            )
            text = result.stdout.strip()

            # If got meaningful text, return it
            if len(text) > 100:
                return {'text': text, 'success': True, 'method': 'native'}
        except Exception:
            pass

        # Fallback to OCR
        if self.use_ocr:
            return self._ocr_pdf(file_path)

        return {'text': '', 'success': False, 'method': 'no_text'}

    def _ocr_pdf(self, file_path: Path) -> Dict:
        """OCR a scanned PDF"""
        try:
            from pdf2image import convert_from_path
            import pytesseract

            images = convert_from_path(str(file_path), dpi=300)
            all_text = []
            confidences = []

            for img in images:
                data = pytesseract.image_to_data(
                    img, lang=self.ocr_languages,
                    output_type=pytesseract.Output.DICT
                )
                text = pytesseract.image_to_string(img, lang=self.ocr_languages)
                all_text.append(text)

                # Get confidence
                confs = [int(c) for c in data['conf'] if c != '-1']
                if confs:
                    confidences.extend(confs)

            avg_conf = sum(confidences) / len(confidences) if confidences else 0

            return {
                'text': '\n\n'.join(all_text),
                'success': True,
                'method': 'ocr',
                'confidence': avg_conf
            }
        except Exception as e:
            return {'text': '', 'success': False, 'method': 'ocr_failed', 'error': str(e)}

    def _extract_word(self, file_path: Path) -> Dict:
        """Extract text from Word document"""
        try:
            if file_path.suffix.lower() == '.docx':
                from docx import Document
                doc = Document(str(file_path))
                text = '\n'.join([p.text for p in doc.paragraphs])
            else:  # .doc
                result = subprocess.run(
                    ['antiword', str(file_path)],
                    capture_output=True, text=True, timeout=30
                )
                text = result.stdout

            return {'text': text.strip(), 'success': True, 'method': 'direct'}
        except Exception as e:
            return {'text': '', 'success': False, 'method': 'error', 'error': str(e)}

    def _extract_excel(self, file_path: Path) -> Dict:
        """Extract text from Excel"""
        try:
            import pandas as pd

            if file_path.suffix.lower() == '.xlsx':
                df = pd.read_excel(str(file_path), sheet_name=None)
            else:
                df = pd.read_excel(str(file_path), sheet_name=None, engine='xlrd')

            all_text = []
            for sheet_name, sheet_df in df.items():
                all_text.append(f"=== Sheet: {sheet_name} ===")
                all_text.append(sheet_df.to_string())

            return {'text': '\n\n'.join(all_text), 'success': True, 'method': 'direct'}
        except Exception as e:
            return {'text': '', 'success': False, 'method': 'error', 'error': str(e)}

    def _extract_image(self, file_path: Path) -> Dict:
        """OCR an image"""
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(str(file_path))

            data = pytesseract.image_to_data(
                img, lang=self.ocr_languages,
                output_type=pytesseract.Output.DICT
            )
            text = pytesseract.image_to_string(img, lang=self.ocr_languages)

            confs = [int(c) for c in data['conf'] if c != '-1']
            avg_conf = sum(confs) / len(confs) if confs else 0

            return {
                'text': text.strip(),
                'success': True,
                'method': 'ocr',
                'confidence': avg_conf
            }
        except Exception as e:
            return {'text': '', 'success': False, 'method': 'error', 'error': str(e)}

    def _extract_text(self, file_path: Path) -> Dict:
        """Extract from plain text files"""
        try:
            encodings = ['utf-8', 'latin-1', 'cp1250', 'iso-8859-2']
            for enc in encodings:
                try:
                    text = file_path.read_text(encoding=enc)
                    return {'text': text.strip(), 'success': True, 'method': 'direct'}
                except UnicodeDecodeError:
                    continue
            return {'text': '', 'success': False, 'method': 'encoding_error'}
        except Exception as e:
            return {'text': '', 'success': False, 'method': 'error', 'error': str(e)}

    def _extract_odf(self, file_path: Path) -> Dict:
        """Extract from OpenDocument files"""
        try:
            from odf import text as odf_text
            from odf.opendocument import load

            doc = load(str(file_path))
            paragraphs = doc.getElementsByType(odf_text.P)
            text = '\n'.join([str(p) for p in paragraphs])

            return {'text': text.strip(), 'success': True, 'method': 'direct'}
        except Exception as e:
            # Fallback: unzip and read content.xml
            try:
                import zipfile
                import xml.etree.ElementTree as ET

                with zipfile.ZipFile(str(file_path)) as zf:
                    content = zf.read('content.xml')
                    root = ET.fromstring(content)
                    text = ' '.join(root.itertext())
                    return {'text': text.strip(), 'success': True, 'method': 'xml'}
            except:
                return {'text': '', 'success': False, 'method': 'error', 'error': str(e)}

    def _extract_ppt(self, file_path: Path) -> Dict:
        """Extract from PowerPoint"""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            all_text = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)

            return {'text': '\n'.join(all_text), 'success': True, 'method': 'direct'}
        except Exception as e:
            return {'text': '', 'success': False, 'method': 'error', 'error': str(e)}


# =============================================================================
# LLM CLASSIFIER
# =============================================================================

class LocalLLMClassifier:
    """Classify documents using local LLM models with consensus

    Backend priority (auto-detected):
    1. MLX (Mac Apple Silicon only) - FASTEST, 2-3x faster than Ollama
    2. Ollama with Metal (Mac) or CUDA (Linux)
    3. Ollama CPU fallback

    Platform-aware configuration:
    - Mac Apple Silicon: Uses MLX if available, else Ollama with Metal
    - Linux NVIDIA (DGX): Can use larger models, vLLM optional
    - Generic: CPU-only, smallest models
    """

    def __init__(self, ollama_url: str = "http://localhost:11434"):
        self.ollama_url = ollama_url
        self.platform = PLATFORM_TYPE
        self.platform_settings = PLATFORM_SETTINGS

        # MLX backend (Mac only - fastest)
        self.mlx_classifier = None
        self.use_mlx = False

        # Try MLX first on Mac (fastest option)
        if MLX_AVAILABLE and PLATFORM_TYPE == "mac_apple_silicon":
            try:
                self.mlx_classifier = MLXClassifier(MLX_MODEL_PATH)
                self.use_mlx = True
                logger.info(f"🚀 Using MLX backend: {MLX_MODEL_PATH.name} (2-3x faster)")
            except Exception as e:
                logger.warning(f"MLX init failed, falling back to Ollama: {e}")

        # Apply platform-specific environment variables
        for key, value in self.platform_settings.get('env_vars', {}).items():
            os.environ.setdefault(key, value)

        self.available_models = self._get_available_models()
        self.selected_models = self._select_models()

        logger.info(f"🖥️  Platform: {self.platform}")
        logger.info(f"🤖 Backend: {'MLX' if self.use_mlx else 'Ollama'}")
        logger.info(f"📊 Models: {['MLX:' + MLX_MODEL_PATH.name] if self.use_mlx else self.selected_models}")
        logger.info(f"⚡ FAST_MODE: {FAST_MODE}")

    def _get_available_models(self) -> List[str]:
        """Get list of available Ollama models"""
        try:
            response = requests.get(f"{self.ollama_url}/api/tags", timeout=5)
            if response.ok:
                return [m['name'] for m in response.json().get('models', [])]
        except:
            pass
        return []

    def _select_models(self) -> List[str]:
        """Select best available models for consensus"""
        selected = []

        # Try preferred models first
        for model in OLLAMA_MODELS:
            if model in self.available_models or model.split(':')[0] in [m.split(':')[0] for m in self.available_models]:
                # Find exact match or closest variant
                for avail in self.available_models:
                    if (model == avail or model.split(':')[0] == avail.split(':')[0]) and avail not in selected:
                        selected.append(avail)
                        break
                if len(selected) >= 3:
                    break

        # Add fallback models if needed
        if len(selected) < 3:
            for model in FALLBACK_MODELS:
                for avail in self.available_models:
                    if model.split(':')[0] == avail.split(':')[0] and avail not in selected:
                        selected.append(avail)
                        break
                if len(selected) >= 3:
                    break

        # If still not enough, use any available models
        if len(selected) < 2:
            for avail in self.available_models:
                if avail not in selected:
                    selected.append(avail)
                    if len(selected) >= 3:
                        break

        return selected[:3]  # Max 3 models

    def _get_fast_model(self) -> Optional[str]:
        """Get best available fast model for hierarchical classification"""
        for model in FAST_MODELS_HIERARCHY:
            base_name = model.split(':')[0]
            for avail in self.available_models:
                if avail == model or avail.split(':')[0] == base_name:
                    return avail
        # Fallback to first selected model
        return self.selected_models[0] if self.selected_models else None

    def _get_big_model(self) -> Optional[str]:
        """Get best available big model for escalation"""
        for model in BIG_MODELS_HIERARCHY:
            base_name = model.split(':')[0]
            for avail in self.available_models:
                if avail == model or avail.split(':')[0] == base_name:
                    return avail
        # No big model available
        return None

    def _call_ollama(self, model: str, prompt: str, timeout: int = 300) -> Optional[str]:
        """Call Ollama model"""
        try:
            response = requests.post(
                f"{self.ollama_url}/api/generate",
                json={
                    "model": model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {"temperature": 0.1, "num_predict": 500}
                },
                timeout=timeout
            )
            if response.ok:
                return response.json().get('response', '').strip()
        except Exception as e:
            logger.warning(f"Ollama call failed for {model}: {e}")
        return None

    def classify(self, text: str, filename: str = "", email_metadata: Dict = None) -> Dict:
        """
        Classify document using hierarchical model orchestration.

        Strategy (HIERARCHICAL_MODE):
        1. Try fast model first (3B quantized - 2x faster)
        2. If confidence < threshold, escalate to big model (14B/32B)
        3. Return combined result with escalation info

        Args:
            text: Extracted document text
            filename: Original filename (for hints)
            email_metadata: Email metadata if from email

        Returns:
            Classification result with votes, consensus, and escalation info
        """
        # Build classification prompt
        prompt = self._build_classification_prompt(text, filename, email_metadata)

        votes = {}
        escalated = False
        fast_model_used = None
        big_model_used = None

        # MLX Backend (Mac Apple Silicon - fastest, 2-3x faster than Ollama)
        if self.use_mlx and self.mlx_classifier:
            try:
                response = self.mlx_classifier.generate(prompt)
                if response:
                    doc_type = self._parse_classification(response)
                    model_name = f"MLX:{MLX_MODEL_PATH.name}" if MLX_MODEL_PATH else "MLX"
                    votes[model_name] = doc_type
                    fast_model_used = model_name
                    logger.debug(f"MLX classification: {doc_type}")
            except Exception as e:
                logger.warning(f"MLX classification failed, falling back to Ollama: {e}")
                self.use_mlx = False  # Disable MLX for this session

        # Ollama Backend with HIERARCHICAL_MODE
        if not votes:
            if HIERARCHICAL_MODE:
                # Step 1: Try fast model first
                fast_model = self._get_fast_model()
                if fast_model:
                    response = self._call_ollama(fast_model, prompt)
                    if response:
                        doc_type = self._parse_classification(response)
                        votes[fast_model] = doc_type
                        fast_model_used = fast_model

                        # Calculate preliminary confidence
                        preliminary_confidence = 1.0 if doc_type != "OTHER" else 0.5

                        # Step 2: Escalate if low confidence or ambiguous result
                        if preliminary_confidence < ESCALATION_THRESHOLD or doc_type == "OTHER":
                            big_model = self._get_big_model()
                            if big_model and big_model != fast_model:
                                logger.info(f"🎭 Escalating to big model: {big_model} (fast: {doc_type}, conf: {preliminary_confidence:.2f})")
                                big_response = self._call_ollama(big_model, prompt)
                                if big_response:
                                    big_doc_type = self._parse_classification(big_response)
                                    votes[big_model] = big_doc_type
                                    big_model_used = big_model
                                    escalated = True
                                    logger.info(f"🎭 Big model result: {big_doc_type}")
            else:
                # Original FAST_MODE: use only first model (3x faster)
                models_to_use = self.selected_models[:1] if FAST_MODE else self.selected_models
                for model in models_to_use:
                    response = self._call_ollama(model, prompt)
                    if response:
                        doc_type = self._parse_classification(response)
                        votes[model] = doc_type

        # Determine consensus
        if not votes:
            return {
                'document_type': 'OTHER',
                'confidence': 0.0,
                'votes': {},
                'consensus': False,
                'escalated': False,
                'fast_model': None,
                'big_model': None
            }

        # Count votes
        type_counts = {}
        for dtype in votes.values():
            type_counts[dtype] = type_counts.get(dtype, 0) + 1

        # Get winner (prefer big model if escalated and disagreement)
        if escalated and len(votes) == 2:
            # If big model was used and disagrees, trust big model
            vote_list = list(votes.values())
            if vote_list[0] != vote_list[1] and big_model_used:
                winner = votes[big_model_used]
                confidence = 0.8  # Lower confidence due to disagreement
            else:
                # Agreement - high confidence
                winner = vote_list[0]
                confidence = 1.0
        else:
            winner = max(type_counts, key=type_counts.get)
            confidence = type_counts[winner] / len(votes)

        consensus = type_counts.get(winner, 0) >= 2

        return {
            'document_type': winner,
            'confidence': confidence,
            'votes': votes,
            'consensus': consensus,
            'escalated': escalated,
            'fast_model': fast_model_used,
            'big_model': big_model_used
        }

    def _build_classification_prompt(self, text: str, filename: str, email_metadata: Dict) -> str:
        """Build prompt for document classification"""
        context = ""
        if filename:
            context += f"Filename: {filename}\n"
        if email_metadata:
            if email_metadata.get('subject'):
                context += f"Email subject: {email_metadata['subject']}\n"
            if email_metadata.get('sender'):
                context += f"Email sender: {email_metadata['sender']}\n"

        text_preview = text[:3000] if len(text) > 3000 else text

        return f"""Analyze this document and classify it into ONE of these types:
{', '.join(DOC_TYPES)}

{context}

Document content:
{text_preview}

Respond with ONLY the document type (one word from the list above), nothing else."""

    def _parse_classification(self, response: str) -> str:
        """Parse classification from LLM response"""
        response_upper = response.upper().strip()

        for dtype in DOC_TYPES:
            if dtype in response_upper:
                return dtype

        return "OTHER"

    def extract_recipient(self, text: str, filename: str = "", email_metadata: Dict = None) -> Optional[RecipientInfo]:
        """Extract recipient information from document"""
        # Priority 1: Email recipient
        if email_metadata and email_metadata.get('recipient'):
            name = email_metadata['recipient']
            rtype = self._detect_recipient_type(name, text)
            return RecipientInfo(name=name, type=rtype, confidence=0.9, source="email")

        # Priority 2: Content analysis with LLM
        if self.selected_models:
            prompt = f"""Analyze this document and identify the primary recipient or addressee.
Look for: "Vážený/á pane/paní", "Dobrý den,", company names, personal names.

Document:
{text[:2000]}

Respond in JSON format:
{{"name": "recipient name", "type": "PERSON_MALE|PERSON_FEMALE|COMPANY|INSTITUTION|UNKNOWN"}}
Only JSON, nothing else."""

            response = self._call_ollama(self.selected_models[0], prompt, timeout=60)
            if response:
                try:
                    # Try to parse JSON
                    import re
                    json_match = re.search(r'\{[^}]+\}', response)
                    if json_match:
                        data = json.loads(json_match.group())
                        if data.get('name') and data.get('type'):
                            return RecipientInfo(
                                name=data['name'],
                                type=data['type'],
                                confidence=0.7,
                                source="content"
                            )
                except:
                    pass

        return None

    def _detect_recipient_type(self, name: str, text: str = "") -> str:
        """Detect recipient type from name and context"""
        name_lower = name.lower()

        # Company indicators
        company_indicators = ['s.r.o.', 'a.s.', 'spol.', 'gmbh', 'inc', 'ltd', 'corp', '@']
        if any(ind in name_lower for ind in company_indicators):
            return "COMPANY"

        # Institution indicators
        institution_indicators = ['ministerstvo', 'úřad', 'soud', 'policie', 'čssz', 'vfp', 'ossz']
        if any(ind in name_lower for ind in institution_indicators):
            return "INSTITUTION"

        # Check for Pan/Paní in text context
        text_lower = text.lower()
        if f"paní {name_lower}" in text_lower or f"pani {name_lower}" in text_lower:
            return "PERSON_FEMALE"
        if f"pan {name_lower}" in text_lower or f"pane {name_lower}" in text_lower:
            return "PERSON_MALE"

        # Female name endings (Czech)
        female_endings = ['ová', 'ová', 'á', 'ice', 'ka']
        if any(name_lower.endswith(e) for e in female_endings):
            return "PERSON_FEMALE"

        return "PERSON_MALE"  # Default assumption

    def suggest_title(self, text: str, doc_type: str, filename: str = "") -> str:
        """Generate suggested title for document"""
        if self.selected_models:
            prompt = f"""Generate a short, descriptive title for this {doc_type} document.
Maximum 60 characters. Use Czech language if the document is in Czech.

Document preview:
{text[:1000]}

Original filename: {filename}

Respond with ONLY the title, nothing else."""

            response = self._call_ollama(self.selected_models[0], prompt, timeout=30)
            if response and len(response) < 100:
                return response.strip('"\'')

        # Fallback: use filename
        if filename:
            return Path(filename).stem[:60]

        return f"{doc_type} document"


# =============================================================================
# UNIFIED DOCUMENT PROCESSOR
# =============================================================================

class UnifiedDocumentProcessor:
    """
    Main class for processing documents from any source

    Workflow:
    1. Receive document (file path + optional source metadata)
    2. Extract text (native or OCR)
    3. Classify with local LLM consensus
    4. Extract recipient
    5. Generate metadata for Paperless-ngx
    """

    def __init__(self,
                 ollama_url: str = "http://localhost:11434",
                 use_ocr: bool = True,
                 progress_file: str = None):
        """
        Initialize processor

        Args:
            ollama_url: URL for Ollama API
            use_ocr: Enable OCR for images and scanned PDFs
            progress_file: Path to JSON file for tracking progress
        """
        self.text_extractor = UnifiedTextExtractor(use_ocr=use_ocr)
        self.classifier = LocalLLMClassifier(ollama_url=ollama_url)
        self.progress_file = Path(progress_file) if progress_file else None

        # Load existing progress
        self.processed_hashes = set()
        if self.progress_file and self.progress_file.exists():
            try:
                data = json.loads(self.progress_file.read_text())
                self.processed_hashes = set(data.get('processed_hashes', []))
            except:
                pass

    def compute_hash(self, file_path: Path) -> str:
        """Compute MD5 hash of file"""
        hash_md5 = hashlib.md5()
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(65536), b''):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def process_document(self,
                        file_path: Union[str, Path],
                        source: Optional[DocumentSource] = None,
                        skip_if_processed: bool = True) -> Optional[ClassificationResult]:
        """
        Process a single document

        Args:
            file_path: Path to document file
            source: Source information (email metadata, etc.)
            skip_if_processed: Skip if already processed

        Returns:
            ClassificationResult or None if skipped/failed
        """
        file_path = Path(file_path)
        start_time = datetime.now()

        # Check file exists and is supported
        if not file_path.exists():
            logger.warning(f"File not found: {file_path}")
            return None

        ext = file_path.suffix.lower()
        if ext not in SUPPORTED_DOCUMENTS:
            logger.debug(f"Unsupported file type: {ext}")
            return None

        # Compute hash and check if already processed
        file_hash = self.compute_hash(file_path)
        if skip_if_processed and file_hash in self.processed_hashes:
            logger.debug(f"Already processed: {file_path.name}")
            return None

        # Create default source if not provided
        if not source:
            source = DocumentSource(
                type="file",
                path=str(file_path),
                original_filename=file_path.name,
                parent_folder=str(file_path.parent)
            )

        # Extract text
        extraction = self.text_extractor.extract(file_path)

        if not extraction.get('success') or not extraction.get('text'):
            logger.warning(f"Text extraction failed: {file_path.name}")
            return None

        text = extraction['text']

        # Build email metadata dict if available
        email_metadata = None
        if source.type == "email":
            email_metadata = {
                'sender': source.email_sender,
                'recipient': source.email_recipient,
                'subject': source.email_subject,
                'date': source.email_date
            }

        # Classify
        classification = self.classifier.classify(
            text=text,
            filename=source.original_filename or file_path.name,
            email_metadata=email_metadata
        )

        # Extract recipient
        recipient = self.classifier.extract_recipient(
            text=text,
            filename=source.original_filename or file_path.name,
            email_metadata=email_metadata
        )

        # Generate title
        title = self.classifier.suggest_title(
            text=text,
            doc_type=classification['document_type'],
            filename=source.original_filename or file_path.name
        )

        # Generate tags
        tags = self._generate_tags(classification['document_type'], recipient, source)

        # Calculate processing time
        processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

        # Create result
        result = ClassificationResult(
            file_path=str(file_path),
            file_hash=file_hash,
            file_type=SUPPORTED_DOCUMENTS[ext],
            source=source,
            document_type=classification['document_type'],
            type_confidence=classification['confidence'],
            type_votes=classification['votes'],
            type_consensus=classification['consensus'],
            recipient=recipient,
            suggested_title=title,
            suggested_tags=tags,
            text_preview=text[:500],
            text_length=len(text),
            processed_at=datetime.now().isoformat(),
            processing_time_ms=processing_time,
            ocr_used=extraction.get('method') == 'ocr',
            models_used=list(classification['votes'].keys())
        )

        # Mark as processed
        self.processed_hashes.add(file_hash)

        return result

    def _generate_tags(self, doc_type: str, recipient: Optional[RecipientInfo], source: DocumentSource) -> List[str]:
        """Generate Paperless-ngx tags"""
        tags = []

        # Document type tag (Czech)
        type_tag_map = {
            "INVOICE": "Faktura",
            "RECEIPT": "Účtenka",
            "CONTRACT": "Smlouva",
            "INSURANCE": "Pojištění",
            "MEDICAL_REPORT": "Lékařská zpráva",
            "BANK_STATEMENT": "Bankovní výpis",
            "TAX_DOCUMENT": "Daňový doklad",
            "LEGAL_DOCUMENT": "Právní dokument",
            "CORRESPONDENCE": "Korespondence",
            "TECHNICAL_MANUAL": "Technická dokumentace",
            "MARKETING": "Marketing",
            "EMAIL": "Email",
            "APPLICATION_FORM": "Žádost",
            "CERTIFICATE": "Certifikát",
        }

        if doc_type in type_tag_map:
            tags.append(type_tag_map[doc_type])

        # Recipient type tag
        if recipient and recipient.type in RECIPIENT_TYPES:
            tags.append(RECIPIENT_TYPES[recipient.type])

        # Source tag
        if source.type == "email":
            tags.append("Z emailu")

        return tags

    def process_directory(self,
                         directory: Union[str, Path],
                         recursive: bool = True,
                         max_files: int = None) -> List[ClassificationResult]:
        """
        Process all documents in a directory

        Args:
            directory: Directory path
            recursive: Include subdirectories
            max_files: Maximum files to process

        Returns:
            List of classification results
        """
        directory = Path(directory)
        results = []
        processed_count = 0

        pattern = "**/*" if recursive else "*"

        for file_path in directory.glob(pattern):
            if max_files and processed_count >= max_files:
                break

            if not file_path.is_file():
                continue

            if file_path.suffix.lower() not in SUPPORTED_DOCUMENTS:
                continue

            result = self.process_document(file_path)
            if result:
                results.append(result)
                processed_count += 1

                if processed_count % 10 == 0:
                    logger.info(f"Processed {processed_count} documents...")

        logger.info(f"Completed: {len(results)} documents processed")
        return results

    def process_email_attachments(self,
                                  attachments_dir: Union[str, Path],
                                  metadata_file: Union[str, Path]) -> List[ClassificationResult]:
        """
        Process email attachments with metadata

        Args:
            attachments_dir: Directory containing extracted attachments
            metadata_file: JSON file with extraction_results.json format

        Returns:
            List of classification results
        """
        attachments_dir = Path(attachments_dir)
        metadata_file = Path(metadata_file)

        # Load metadata
        if not metadata_file.exists():
            logger.error(f"Metadata file not found: {metadata_file}")
            return []

        metadata = json.loads(metadata_file.read_text())
        files_metadata = {f['saved_as']: f for f in metadata.get('extracted_files', [])}

        results = []

        for file_path in attachments_dir.iterdir():
            if not file_path.is_file():
                continue

            # Get metadata for this file
            file_meta = files_metadata.get(str(file_path), {})

            # Create source from metadata
            source = DocumentSource(
                type="email",
                path=str(file_path),
                email_sender=file_meta.get('email_from'),
                email_recipient=None,  # Could be extracted
                email_subject=file_meta.get('email_subject'),
                email_date=file_meta.get('email_date'),
                original_filename=file_meta.get('original_filename'),
                parent_folder=str(attachments_dir)
            )

            result = self.process_document(file_path, source=source)
            if result:
                results.append(result)

        return results

    def save_results(self, results: List[ClassificationResult], output_file: Union[str, Path]):
        """Save results to JSON file"""
        output_file = Path(output_file)

        data = {
            'timestamp': datetime.now().isoformat(),
            'total_processed': len(results),
            'stats': self._compute_stats(results),
            'results': [asdict(r) for r in results]
        }

        output_file.write_text(json.dumps(data, indent=2, ensure_ascii=False))
        logger.info(f"Results saved to: {output_file}")

    def _compute_stats(self, results: List[ClassificationResult]) -> Dict:
        """Compute statistics from results"""
        stats = {
            'by_type': {},
            'by_recipient_type': {},
            'consensus_rate': 0,
            'ocr_used_count': 0,
            'avg_processing_time_ms': 0
        }

        if not results:
            return stats

        for r in results:
            # By document type
            stats['by_type'][r.document_type] = stats['by_type'].get(r.document_type, 0) + 1

            # By recipient type
            if r.recipient:
                rtype = r.recipient.type
                stats['by_recipient_type'][rtype] = stats['by_recipient_type'].get(rtype, 0) + 1

            # OCR usage
            if r.ocr_used:
                stats['ocr_used_count'] += 1

        # Consensus rate
        consensus_count = sum(1 for r in results if r.type_consensus)
        stats['consensus_rate'] = consensus_count / len(results) if results else 0

        # Average processing time
        total_time = sum(r.processing_time_ms for r in results)
        stats['avg_processing_time_ms'] = total_time / len(results) if results else 0

        return stats

    def save_progress(self):
        """Save progress to file"""
        if self.progress_file:
            data = {
                'timestamp': datetime.now().isoformat(),
                'processed_hashes': list(self.processed_hashes)
            }
            self.progress_file.write_text(json.dumps(data, indent=2))


# =============================================================================
# CLI
# =============================================================================

def main():
    import argparse

    parser = argparse.ArgumentParser(description='Unified Document Processor')
    parser.add_argument('--source', required=True, help='Source directory or file')
    parser.add_argument('--output', required=True, help='Output JSON file')
    parser.add_argument('--email-metadata', help='Email extraction metadata JSON')
    parser.add_argument('--ollama-url', default='http://localhost:11434', help='Ollama API URL')
    parser.add_argument('--no-ocr', action='store_true', help='Disable OCR')
    parser.add_argument('--max-files', type=int, help='Maximum files to process')
    parser.add_argument('--progress-file', help='Progress tracking file')

    args = parser.parse_args()

    processor = UnifiedDocumentProcessor(
        ollama_url=args.ollama_url,
        use_ocr=not args.no_ocr,
        progress_file=args.progress_file
    )

    source_path = Path(args.source)

    if args.email_metadata:
        # Process email attachments
        results = processor.process_email_attachments(
            attachments_dir=source_path,
            metadata_file=args.email_metadata
        )
    elif source_path.is_file():
        # Process single file
        result = processor.process_document(source_path)
        results = [result] if result else []
    else:
        # Process directory
        results = processor.process_directory(
            directory=source_path,
            max_files=args.max_files
        )

    # Save results
    processor.save_results(results, args.output)
    processor.save_progress()

    print(f"\nProcessed: {len(results)} documents")
    print(f"Results saved to: {args.output}")


if __name__ == "__main__":
    main()
